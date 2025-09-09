import io, math, ast, os, json, tempfile
import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
from scipy import stats as sps
import statsmodels.api as sm
from statsmodels.stats.diagnostic import het_breuschpagan
from statsmodels.graphics.gofplots import qqplot as sm_qqplot

from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from PIL import Image as PILImage

from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PPTInches
import zipfile

st.set_page_config(page_title="FrostNode | Statistics Toolset", layout="wide")

# --------------------------- Knowledge: Definitions & Tips ---------------------------
KB = {
"Unit 1 — Introduction": {
"definitions": [
("Statistics","Collecting, organizing, analyzing, and interpreting data to make decisions."),
("Population vs Sample","Entire group vs. subset observed."),
("Parameter vs Statistic","Describes population vs. describes sample."),
("Observational vs Experiment","Observe only vs. assign treatments (supports causal claims if well-designed).")
],
"tips":[
"Define the research question first.",
"Random selection reduces bias; random assignment enables causal inference."
]
},
"Unit 2–4 — Descriptives": {
"definitions":[
("Mean/Median/Mode","Center measures; median resistant to outliers."),
("Std Dev (s) & IQR","Spread measures; IQR = Q3−Q1 (resistant)."),
("Five-number summary","Min, Q1, Median, Q3, Max (for boxplots).")
],
"tips":[
"Skewed/outliers → use median & IQR.",
"Symmetric → mean & SD are fine."
]
},
"Unit 5–6 — Correlation & Regression":{
"definitions":[
("Correlation (r)","Strength/direction of linear association (−1 to +1)."),
("Regression line","Best-fit line minimizing squared residuals."),
("Residual","Observed − Predicted; should bounce around 0."),
("r²","Share of variability in y explained by the linear model.")
],
"tips":[
"Correlation ≠ causation.",
"Check residual plot & avoid extrapolation."
]
},
"Unit 9–10 — CIs & Probability":{
"definitions":[
("Confidence Interval","Range of plausible values for a parameter."),
("Sampling distribution","Distribution of a statistic over many samples.")
],
"tips":[
"Higher confidence → wider interval.",
"For proportion CIs, check np̂ and n(1−p̂) ≥ ~10."
]
}
}

def render_kb():
    with st.sidebar.expander("📚 Definitions & Tips (quick help)", expanded=False):
        for sec,items in KB.items():
            st.markdown(f"**{sec}**")
            st.write("Definitions:")
            for term,defn in items["definitions"]:
                st.write(f"- **{term}** — {defn}")
            st.write("Tips:")
            for tip in items["tips"]:
                st.write(f"- {tip}")
            st.write("---")

# --------------------------- Utilities ---------------------------
def sturges_k(n):
    return max(1, int(round(1 + np.log2(max(n,1)))))

def fig_bytes(fig, dpi=200):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight")
    buf.seek(0)
    return buf.read()

def safe_eval(expr, var='x'):
    allowed_nodes = (ast.Expression, ast.BinOp, ast.UnaryOp, ast.Num, ast.Load,
                     ast.Add, ast.Sub, ast.Mult, ast.Div, ast.Pow, ast.USub,
                     ast.Call, ast.Name)
    allowed_names = {var} | {'sin','cos','tan','arcsin','arccos','arctan','exp','log','log10','sqrt','abs','pi','e'}
    node = ast.parse(expr, mode='eval')
    for n in ast.walk(node):
        if not isinstance(n, allowed_nodes):
            raise ValueError(f"Unsupported expression element: {type(n).__name__}")
        if isinstance(n, ast.Call):
            if not isinstance(n.func, ast.Name) or n.func.id not in allowed_names:
                raise ValueError("Only basic math functions are allowed.")
        if isinstance(n, ast.Name) and n.id not in allowed_names:
            raise ValueError(f"Unknown name: {n.id}")
    code = compile(node, "<expr>", "eval")
    def f(x):
        env = {var: x,'sin':np.sin,'cos':np.cos,'tan':np.tan,'arcsin':np.arcsin,'arccos':np.arccos,'arctan':np.arctan,
               'exp':np.exp,'log':np.log,'log10':np.log10,'sqrt':np.sqrt,'abs':np.abs,'pi':np.pi,'e':np.e}
        return eval(code, {"__builtins__":{}}, env)
    return f

def to_numeric_series(s, drop_log):
    s2 = pd.to_numeric(s, errors="coerce")
    dropped = int(s2.isna().sum()) - int(s.isna().sum())
    if dropped>0:
        drop_log.append(f"Dropped {dropped} non-numeric cell(s) after coercion in column '{s.name}'.")
    return s2

def skew_narrative(series):
    s = pd.Series(series)
    if s.size<3: return "Not enough data for shape notes."
    mean, median = s.mean(), s.median()
    if s.std(ddof=1) == 0: return "Data show no variation."
    if mean > median*1.05: return "Right-skewed; consider median & IQR."
    if mean < median*0.95: return "Left-skewed; consider median & IQR."
    return "Roughly symmetric; mean & SD are fine."

def fisher_r_ci(r, n, conf=0.95):
    # Need at least 4 points and a finite r
    if n is None or n < 4 or not np.isfinite(r):
        return (np.nan, np.nan)
    # Clip r strictly inside (-1, 1) to avoid division-by-zero in Fisher z
    r = float(np.clip(r, -0.999999, 0.999999))
    z = 0.5*np.log((1+r)/(1-r))
    se = 1.0/np.sqrt(n - 3.0)
    zcrit = sps.norm.ppf(0.5 + conf/2.0)
    lo = z - zcrit*se
    hi = z + zcrit*se
    rlo = (np.exp(2*lo)-1)/(np.exp(2*lo)+1)
    rhi = (np.exp(2*hi)-1)/(np.exp(2*hi)+1)
    return (rlo, rhi)

# --------------------------- Import helpers ---------------------------
with st.sidebar:
    st.header("Import data")
    up = st.file_uploader("CSV or Excel (.xlsx)", type=["csv","xlsx"])
    sample = st.selectbox("Or load a sample", ["—","Hours_vs_Scores","RN_Salaries"], index=0)
    st.caption("CSV helpers")
    delim = st.text_input("Delimiter (CSV only)", value=",")
    has_header = st.checkbox("First row is header", value=True)
    coerce = st.checkbox("Coerce to numeric (drop non-numeric)", value=True)

render_kb()

df = None
drop_log = []
sheet = None

if up is not None:
    if up.name.lower().endswith(".csv"):
        df = pd.read_csv(up, delimiter=delim, header=0 if has_header else None)
        if not has_header:
            df.columns = [f"col_{i+1}" for i in range(df.shape[1])]
    else:
        xls = pd.ExcelFile(up)
        sheet = st.sidebar.selectbox("Select sheet", xls.sheet_names, index=0)
        df = xls.parse(sheet)

if df is None and sample != "—":
    if sample=="Hours_vs_Scores":
        df = pd.DataFrame({"Hours":[0,2,4,5,5,5,6,7,8],
                           "Score":[40,51,64,69,73,75,93,90,95]})
    elif sample=="RN_Salaries":
        df = pd.DataFrame({"Years":[0.5,2,4,5,7,9,10,12.5,13,16,18,20,25],
                           "Salary_k":[48.3,53.4,58.5,63.4,65.7,67.3,69.8,71.8,73.2,75.5,74.3,78.9,76.6]})

if df is None:
    st.title("📈 FrostNode | Statistics Toolset")
    st.info("Upload a CSV/Excel file or choose a sample in the sidebar. Use the helpers to adjust delimiter and headers.")
    st.stop()

# Coerce numeric columns if requested
if coerce:
    for c in df.columns:
        before_na = df[c].isna().sum()
        df[c] = pd.to_numeric(df[c], errors="ignore")
        if not pd.api.types.is_numeric_dtype(df[c]):
            s = pd.to_numeric(df[c], errors="coerce")
            dropped = int(s.isna().sum()) - int(df[c].isna().sum())
            if dropped>0:
                drop_log.append(f"Dropped {dropped} non-numeric cell(s) in column '{c}'.")
            df[c] = s

st.title("📈 FrostNode | Statistics Toolset")
st.subheader("Data preview")
st.dataframe(df.head(25), use_container_width=True)
if drop_log:
    st.caption("Import log:")
    for line in drop_log:
        st.write("• " + line)

num_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
if not num_cols:
    st.error("No numeric columns detected after import.")
    st.stop()

# Session controls
with st.sidebar:
    st.header("Session")
    # Save session selections
    if st.button("Save session (.json)"):
        sess = {"columns": list(df.columns), "num_cols": num_cols}
        st.download_button("Download session.json", data=json.dumps(sess).encode("utf-8"),
                           file_name="frostnode_session.json", mime="application/json", key="sessdl")
    state_up = st.file_uploader("Load session (.json)", type=["json"], key="stateup")
    if state_up is not None:
        st.success("Session file received (columns read).")

# Mode & analysis setup
with st.sidebar:
    mode = st.radio("Mode", ["Simple statistics","Advanced statistics"], index=0)
    alpha = st.select_slider("α (significance)", [0.10,0.05,0.02,0.01,0.001], value=0.05)
    xcol = st.selectbox("X (numeric)", num_cols, index=0)
    ycol = st.selectbox("Y (numeric, optional)", ["—"] + num_cols, index=(1 if len(num_cols)>1 else 0))
    x0_input = st.text_input("Prediction x₀ (optional)", value="")

X = df[xcol].dropna().astype(float).values
Y = None if ycol=="—" else df[ycol].dropna().astype(float).values

# --------------------------- Figures collector for ZIP/exports ---------------------------
assets = {"tables":{}, "figs":{}, "texts":{}}

def add_png(name, fig):
    assets["figs"][f"{name}.png"] = fig_bytes(fig)

def add_csv(name, df_):
    assets["tables"][f"{name}.csv"] = df_.to_csv(index=False).encode("utf-8")

def add_text(name, s):
    assets["texts"][f"{name}.txt"] = s.encode("utf-8")

# --------------------------- Simple Mode ---------------------------
def hist_plot(data, bins, xlab, title):
    fig, ax = plt.subplots(figsize=(6,4))
    ax.hist(data, bins=bins, edgecolor="white")
    ax.set_xlabel(xlab); ax.set_ylabel("Count"); ax.set_title(title)
    ax.grid(alpha=0.25)
    return fig

def box_plot(data, xlab, title):
    fig, ax = plt.subplots(figsize=(4,4))
    ax.boxplot(data, vert=True, patch_artist=True)
    ax.set_xticklabels([xlab])
    ax.set_title(title)
    ax.grid(alpha=0.25)
    return fig

def regression_xy(x, y):
    x = pd.Series(x).astype(float).to_numpy()
    y = pd.Series(y).astype(float).to_numpy()
    n = min(len(x), len(y))
    x, y = x[:n], y[:n]
    r_p, p_r = sps.pearsonr(x, y)
    Xex = sm.add_constant(x)
    model = sm.OLS(y, Xex).fit()
    m = float(model.params[1]); b = float(model.params[0])
    r2 = float(model.rsquared)
    resid = y - model.fittedvalues
    sse = float(np.sum(resid**2))
    return dict(n=n, r=float(r_p), p_r=float(p_r), m=m, b=b, r2=r2, sse=sse, model=model, resid=resid, x=x, y=y)

def simple_scatter_plot(x, y, m, b, xlab, ylab, title):
    fig, ax = plt.subplots(figsize=(6,4))
    ax.scatter(x, y, alpha=0.85, edgecolor="white")
    xx = np.linspace(np.min(x), np.max(x), 100)
    ax.plot(xx, m*xx + b, linewidth=2)
    ax.set_xlabel(xlab); ax.set_ylabel(ylab); ax.set_title(title)
    ax.grid(alpha=0.25)
    return fig

def residual_plot(x, resid, xlab, title):
    fig, ax = plt.subplots(figsize=(6,4))
    ax.scatter(x, resid, alpha=0.85, edgecolor="white")
    ax.axhline(0, color="red", linestyle="--", linewidth=1)
    ax.set_xlabel(xlab); ax.set_ylabel("Residual (y − ŷ)"); ax.set_title(title)
    ax.grid(alpha=0.25)
    return fig

def qq_plot(resid, title):
    fig = sm_qqplot(resid, line='45', fit=True)
    fig.suptitle(title)
    return fig

# Transform toggles
transform_choice = "None"
if mode == "Advanced statistics" and (Y is not None):
    transform_choice = st.selectbox("Transform data for regression", ["None","log(y)","log(x)","log(x), log(y)"], index=0)
    def transform_xy(x, y, choice):
        tx, ty = x.copy(), y.copy()
        if "log(x)" in choice:
            if np.any(tx<=0): 
                st.warning("log(x) selected but x has non-positive values. Those rows will be dropped for the transformed fit.")
            mask = tx>0
            tx, ty = tx[mask], ty[mask]
            tx = np.log(tx)
        if "log(y)" in choice:
            if np.any(ty<=0):
                st.warning("log(y) selected but y has non-positive values. Those rows will be dropped for the transformed fit.")
            mask = ty>0
            tx, ty = tx[mask], ty[mask]
            ty = np.log(ty)
        return tx, ty
else:
    transform_choice = "None"

# --------------------------- UI: Simple or Advanced ---------------------------
if mode == "Simple statistics":
    st.header("Simple tools")
    # Quick stats
    s = pd.Series(X)
    q = s.quantile([0,0.25,0.5,0.75,1.0])
    simple_tbl = pd.DataFrame({
        "Metric":["Count","Mean","Median","Std (sample)","Min","Q1","Q2/Median","Q3","Max","IQR"],
        "Value":[s.size, s.mean(), s.median(), s.std(ddof=1),
                 q.loc[0.0], q.loc[0.25], q.loc[0.5], q.loc[0.75], q.loc[1.0],
                 q.loc[0.75]-q.loc[0.25]]
    })
    st.subheader("Quick Stats (X)")
    st.dataframe(simple_tbl, use_container_width=True)
    add_csv("quick_stats_X", simple_tbl)
    # Histogram
    bins = st.number_input("Bins (Sturges default)", value=int(sturges_k(len(X))), min_value=1)
    figH = hist_plot(X, int(bins), xlab=xcol, title=f"Histogram of {xcol}")
    st.pyplot(figH, use_container_width=True)
    add_png(f"hist_{xcol}", figH)
    st.caption("Narrative: " + skew_narrative(X))
    add_text("histogram_narrative", f"{xcol}: " + skew_narrative(X))
    st.download_button("Download histogram PNG", data=assets["figs"][f"hist_{xcol}.png"], file_name=f"hist_{xcol}.png", mime="image/png")
    # Scatter + line
    if Y is not None:
        res = regression_xy(X, Y)
        st.divider()
        st.subheader("Quick Scatter + Line (if Y selected)")
        st.write(f"**Regression line**: ŷ = {res['m']:.6g}·x + {res['b']:.6g} | r={res['r']:.4f}, r²={res['r2']:.4f}")
        st.write(f"Correlation test (H₀: ρ=0) p = {res['p_r']:.4g} (α={alpha}) → " +
                 ("**significant**" if res['p_r']<alpha else "**not significant**"))
        figS = simple_scatter_plot(res['x'], res['y'], res['m'], res['b'], xlab=xcol, ylab=ycol, title=f"{ycol} vs {xcol}")
        st.pyplot(figS, use_container_width=True)
        add_png(f"scatter_{ycol}_vs_{xcol}", figS)
        st.download_button("Download scatter+line PNG", data=assets["figs"][f"scatter_{ycol}_vs_{xcol}.png"],
                           file_name=f"scatter_{ycol}_vs_{xcol}.png", mime="image/png")
        conclusion = ("There is " + ("statistically significant " if res['p_r']<alpha else "no statistically significant ") +
                      f"linear association between {xcol} and {ycol} at α={alpha}. Fitted line: ŷ={res['m']:.3g}x+{res['b']:.3g}; r²≈{res['r2']:.2f}.")
        st.caption("Conclusion: " + conclusion)
        add_text("regression_conclusion", conclusion)

else:
    st.header("Advanced tools")

    tab1, tab2, tab3, tab4, tab5 = st.tabs(["Descriptive & Graphs","Correlation & Regression","Confidence Intervals","Function Plotter","Exports"])

    with tab1:
        st.subheader("Descriptive statistics (X)")
        desc = pd.Series(X).describe(percentiles=[0.25,0.5,0.75])
        extra = pd.Series({"IQR": pd.Series(X).quantile(0.75) - pd.Series(X).quantile(0.25)})
        dfx = pd.concat([desc, extra]).to_frame("Value").reset_index().rename(columns={"index":"Metric"})
        st.dataframe(dfx, use_container_width=True)
        add_csv("desc_X", dfx)
        st.caption("Narrative: " + skew_narrative(X))
        add_text("desc_narrative", skew_narrative(X))
        # Frequency
        k = st.number_input("Bins (k)", value=int(sturges_k(len(X))), min_value=1)
        counts, bins = np.histogram(X, bins=int(k))
        freq_df = pd.DataFrame({
            "Bin":[f"[{bins[i]:.6g},{bins[i+1]:.6g}{']' if i==len(counts)-1 else ')'}" for i in range(len(counts))],
            "Count":counts
        })
        st.subheader("Frequency table (Sturges)")
        st.dataframe(freq_df, use_container_width=True)
        add_csv("frequency_X", freq_df)
        # Plots
        cA, cB = st.columns(2)
        with cA:
            st.subheader("Histogram")
            figH = hist_plot(X, int(k), xlab=xcol, title=f"Histogram of {xcol}")
            st.pyplot(figH, use_container_width=True)
            add_png(f"adv_hist_{xcol}", figH)
            st.download_button("Download PNG", data=assets["figs"][f"adv_hist_{xcol}.png"], file_name=f"adv_hist_{xcol}.png", mime="image/png")
        with cB:
            st.subheader("Boxplot")
            figB = box_plot(X, xlab=xcol, title=f"Boxplot of {xcol}")
            st.pyplot(figB, use_container_width=True)
            add_png(f"adv_box_{xcol}", figB)
            st.download_button("Download PNG", data=assets["figs"][f"adv_box_{xcol}.png"], file_name=f"adv_box_{xcol}.png", mime="image/png")

    with tab2:
        if Y is None:
            st.info("Select a Y column to enable correlation & regression.")
        else:
            # Transform and fit
            tx, ty = (X.copy(), Y.copy())
            if transform_choice != "None":
                tx, ty = (lambda a,b: (a,b))(*transform_xy(X, Y, transform_choice))
            res = regression_xy(tx, ty)

            st.subheader("Key metrics")
            c1, c2, c3, c4, c5 = st.columns(5)
            c1.metric("r (Pearson)", f"{res['r']:.4f}")
            c2.metric("r²", f"{res['r2']:.4f}")
            c3.metric("SSE", f"{res['sse']:.4g}")
            c4.metric("n", f"{res['n']}")
            # Spearman
            rho, p_spear = sps.spearmanr(tx, ty)
            c5.metric("Spearman ρ", f"{rho:.4f}")
            st.caption(f"Spearman test p = {p_spear:.4g}")

            st.write(f"**Regression line**: ŷ = {res['m']:.6g}·x + {res['b']:.6g}")
            if transform_choice!="None":
                st.info(f"Transformed fit used: {transform_choice}. If log(y) used, back-transform prediction via exp(ŷ).")

            # CIs for slope/intercept
            conf = res["model"].conf_int(alpha=alpha)
            b0_lo, b0_hi = conf[0]
            b1_lo, b1_hi = conf[1]
            st.write(f"{int((1-alpha)*100)}% CIs — Intercept b₀: ({b0_lo:.4g}, {b0_hi:.4g}), Slope b₁: ({b1_lo:.4g}, {b1_hi:.4g})")

            # CI for r (Fisher z)
            rlo, rhi = fisher_r_ci(res["r"], res["n"], conf=1-alpha)
            if np.isfinite(rlo):
                st.write(f"{int((1-alpha)*100)}% CI for r: ({rlo:.4f}, {rhi:.4f})")

            # Plots
            sc1, sc2 = st.columns(2)
            with sc1:
                st.subheader("Scatter + fitted line")
                xlab_str = (f"log({xcol})" if "log(x)" in transform_choice else xcol)
                ylab_str = (f"log({ycol})" if "log(y)" in transform_choice else ycol)
                figS = simple_scatter_plot(res['x'], res['y'], res['m'], res['b'], xlab=xlab_str, ylab=ylab_str, title=f"{ylab_str} vs {xlab_str}")

                st.pyplot(figS, use_container_width=True)
                add_png(f"adv_scatter_{ycol}_vs_{xcol}", figS)
                st.download_button("Download PNG", data=assets["figs"][f"adv_scatter_{ycol}_vs_{xcol}.png"], file_name=f"adv_scatter_{ycol}_vs_{xcol}.png", mime="image/png")
            with sc2:
                st.subheader("Residual plot")
                figR = residual_plot(res['x'], res['resid'], xlab=(f"log({xcol})" if "log(x)" in transform_choice else xcol), title="Residuals vs X")
                st.pyplot(figR, use_container_width=True)
                add_png(f"adv_residuals_{ycol}_vs_{xcol}", figR)
                st.download_button("Download PNG", data=assets["figs"][f"adv_residuals_{ycol}_vs_{xcol}.png"], file_name=f"adv_residuals_{ycol}_vs_{xcol}.png", mime="image/png")
                st.caption("Look for random bounce around 0; patterns suggest nonlinearity or changing variance.")

            # Assumption checks
            st.subheader("Assumption checks")
            # Q-Q plot
            figQ = qq_plot(res["resid"], "Residuals Q-Q plot")
            st.pyplot(figQ, use_container_width=True)
            add_png("qq_residuals", figQ)
            # Normality test
            if res["n"]>=3:
                try:
                    shW, pW = sps.shapiro(res["resid"])
                    st.write(f"Shapiro–Wilk p = {pW:.4g} (p>α suggests residuals are approximately normal)")
                except Exception:
                    st.caption("Shapiro–Wilk not computed (large n).")
            # Heteroskedasticity hint (Breusch-Pagan)
            try:
                lm_stat, lm_p, f_stat, f_p = het_breuschpagan(res["resid"], sm.add_constant(res["x"]))
                st.write(f"Breusch–Pagan p = {lm_p:.4g} (small p suggests non-constant variance)")
            except Exception:
                pass

            # Prediction & intervals
            st.subheader("Prediction & intervals")
            x0 = st.text_input("x₀ (single value)", value=x0_input or "")
            if x0.strip():
                try:
                    x0v = float(x0)
                    new = pd.DataFrame({"const":[1.0], res["model"].model.exog_names[1]:[x0v]})
                    pred = res["model"].get_prediction(new).summary_frame(alpha=alpha)
                    yhat = float(pred["mean"].iloc[0])
                    loP, hiP = float(pred["obs_ci_lower"].iloc[0]), float(pred["obs_ci_upper"].iloc[0])
                    loM, hiM = float(pred["mean_ci_lower"].iloc[0]), float(pred["mean_ci_upper"].iloc[0])
                    if "log(y)" in transform_choice:
                        st.write(f"Transformed ŷ (log-scale) = {yhat:.4g}; back-transform exp(ŷ) ≈ {math.exp(yhat):.4g}")
                    st.write(f"ŷ({x0v}) = **{yhat:.4g}**")
                    st.write(f"{int((1-alpha)*100)}% Prediction interval: ({loP:.4g}, {hiP:.4g})")
                    st.write(f"{int((1-alpha)*100)}% Mean response interval: ({loM:.4g}, {hiM:.4g})")
                    xmin, xmax = float(np.min(res["x"])), float(np.max(res["x"]))
                    if x0v < xmin or x0v > xmax:
                        st.warning(f"x₀ is outside observed x-range [{xmin:.4g}, {xmax:.4g}] — extrapolation.")
                    add_text("prediction_report", f"x0={x0v}, yhat={yhat}, PI=({loP},{hiP}), MRI=({loM},{hiM})")
                except Exception as e:
                    st.error(f"Could not compute prediction: {e}")
            else:
                st.info("Enter an x₀ to compute prediction and intervals.")

            with st.expander("Model summary (statsmodels)"):
                st.text(res["model"].summary())

            # Export regression CSV
            out = pd.DataFrame({"x": res["x"], "y": res["y"], "yhat": res["model"].fittedvalues, "residual": res["resid"]})
            st.download_button("Download regression CSV", data=out.to_csv(index=False).encode("utf-8"),
                               file_name=f"regression_{ycol}_on_{xcol}.csv", mime="text/csv")
            add_csv(f"regression_{ycol}_on_{xcol}", out)

    with tab3:
        st.subheader("Confidence intervals")
        ciA, ciB = st.columns(2)
        with ciA:
            st.markdown("**Mean (z/t)**")
            xbar = st.number_input("x̄", value=float(np.nanmean(X)))
            sd = st.number_input("s or σ", value=float(np.nanstd(X, ddof=1)))
            n = st.number_input("n", min_value=1, value=int(len(X)))
            use_t = st.checkbox("Use t (σ unknown; small n)", value=True)
            if st.button("Compute CI (mean)"):
                if use_t:
                    tcrit = sps.t.ppf(1-alpha/2, df=n-1)
                    moe = tcrit*sd/math.sqrt(n)
                    st.success(f"{int((1-alpha)*100)}% CI: ({xbar-moe:.6g}, {xbar+moe:.6g}) using t*={tcrit:.3f}")
                    add_text("CI_mean", f"Mean CI: ({xbar-moe}, {xbar+moe}), t*={tcrit}")
                else:
                    zcrit = sps.norm.ppf(1-alpha/2)
                    moe = zcrit*sd/math.sqrt(n)
                    st.success(f"{int((1-alpha)*100)}% CI: ({xbar-moe:.6g}, {xbar+moe:.6g}) using z*={zcrit:.3f}")
                    add_text("CI_mean", f"Mean CI: ({xbar-moe}, {xbar+moe}), z*={zcrit}")
            st.caption("Narrative: Interval shows plausible values for the population mean under model assumptions.")

        with ciB:
            st.markdown("**Proportion**")
            x_succ = st.number_input("x (successes)", value=60)
            n_prop = st.number_input("n (trials)", value=100)
            if st.button("Compute CI (proportion)"):
                p_hat = x_succ/n_prop if n_prop>0 else np.nan
                zcrit = sps.norm.ppf(1-alpha/2)
                se = math.sqrt(p_hat*(1-p_hat)/n_prop)
                moe = zcrit*se
                st.success(f"p̂={p_hat:.4f}; {int((1-alpha)*100)}% CI: ({p_hat-moe:.6g}, {p_hat+moe:.6g}) with z*={zcrit:.3f}")
                add_text("CI_proportion", f"Prop CI: ({p_hat-moe}, {p_hat+moe}), z*={zcrit}")
                if n_prop*p_hat < 10 or n_prop*(1-p_hat) < 10:
                    st.warning("np̂ and n(1−p̂) should both be ≥ ~10 for normal-approx CI.")

        # Excel Formula Generator (with absolute range option)
        st.divider()
        st.subheader("Excel Formula Generator")
        abs_refs = st.checkbox("Use absolute ranges ($A$1:$A$101)", value=True)
        Y_rng = st.text_input("Excel Y range", value="$B$2:$B$101" if abs_refs else "B2:B101")
        X_rng = st.text_input("Excel X range", value="$A$2:$A$101" if abs_refs else "A2:A101")
        x0_cell = st.text_input("Cell for x₀", value="$E$2" if abs_refs else "E2")
        formulas = {
            "Correlation (r)": f"=CORREL({Y_rng}, {X_rng})",
            "R-squared": f"=RSQ({Y_rng}, {X_rng})",
            "Slope (m)": f"=SLOPE({Y_rng}, {X_rng})",
            "Intercept (b)": f"=INTERCEPT({Y_rng}, {X_rng})",
            "Regression (ŷ at x)": f"=FORECAST.LINEAR({x0_cell}, {Y_rng}, {X_rng})",
            "Std error (STEYX)": f"=STEYX({Y_rng}, {X_rng})",
            "Mean (Y)": f"=AVERAGE({Y_rng})",
            "Median (Y)": f"=MEDIAN({Y_rng})",
            "Mode (Y)": f"=MODE.SNGL({Y_rng})",
            "Std Dev (Y sample)": f"=STDEV.S({Y_rng})",
            "Variance (Y sample)": f"=VAR.S({Y_rng})",
            "Q1/Q2/Q3": f"=QUARTILE.EXC({Y_rng},1/2/3)",
            "Percentile (90%)": f"=PERCENTILE.INC({Y_rng},0.9)",
            "Normal CDF": "=NORM.S.DIST(z, TRUE)",
            "Normal INV": "=NORM.S.INV(p)"
        }
        st.json(formulas)
        st.download_button("Download Formula Sheet (.txt)", data="\n".join([f"{k}: {v}" for k,v in formulas.items()]).encode("utf-8"),
                           file_name="Excel_Formulas.txt", mime="text/plain")

    with tab4:
        st.subheader("Function plotter (make & download charts)")
        expr = st.text_input("Enter f(x)", value="2*x + 3")
        dom1, dom2 = st.columns(2)
        with dom1:
            a = st.number_input("Domain start (a)", value=-10.0)
        with dom2:
            b = st.number_input("Domain end (b)", value=10.0)
        N = st.number_input("Number of points", min_value=10, max_value=10000, value=200)
        go = st.button("Plot f(x)")
        if go:
            try:
                def safe_eval(expr, var='x'):
                    allowed_nodes = (ast.Expression, ast.BinOp, ast.UnaryOp, ast.Num, ast.Load, ast.Add, ast.Sub, ast.Mult, ast.Div, ast.Pow, ast.USub, ast.Call, ast.Name)
                    allowed_names = {var} | {'sin','cos','tan','arcsin','arccos','arctan','exp','log','log10','sqrt','abs','pi','e'}
                    node = ast.parse(expr, mode='eval')
                    for n in ast.walk(node):
                        if not isinstance(n, allowed_nodes): raise ValueError(f"Unsupported element: {type(n).__name__}")
                        if isinstance(n, ast.Call):
                            if not isinstance(n.func, ast.Name) or n.func.id not in allowed_names: raise ValueError("Only basic math functions are allowed.")
                        if isinstance(n, ast.Name) and n.id not in allowed_names: raise ValueError(f"Unknown name: {n.id}")
                    code = compile(node, "<expr>", "eval")
                    def f(x):
                        env = {var: x,'sin':np.sin,'cos':np.cos,'tan':np.tan,'arcsin':np.arcsin,'arccos':np.arccos,'arctan':np.arctan,'exp':np.exp,'log':np.log,'log10':np.log10,'sqrt':np.sqrt,'abs':np.abs,'pi':np.pi,'e':np.e}
                        return eval(code, {"__builtins__":{}}, env)
                    return f
                f = safe_eval(expr, 'x')
                xs = np.linspace(a, b, int(N))
                ys = f(xs)
                fig, ax = plt.subplots(figsize=(6,4))
                ax.plot(xs, ys, linewidth=2)
                ax.set_title(f"f(x) = {expr}")
                ax.set_xlabel("x"); ax.set_ylabel("f(x)"); ax.grid(alpha=0.25)
                st.pyplot(fig, use_container_width=True)
                add_png("function_plot", fig)
                st.download_button("Download plot PNG", data=assets["figs"]["function_plot.png"], file_name="function_plot.png", mime="image/png")
                data = pd.DataFrame({"x": xs, "f(x)": ys})
                st.dataframe(data.head(20), use_container_width=True)
                add_csv("function_values", data)
                st.download_button("Download table CSV", data=assets["tables"]["function_values.csv"], file_name="function_values.csv", mime="text/csv")
            except Exception as e:
                st.error(f"Error: {e}")

    with tab5:
        st.subheader("One-click exports")
        # Download-All ZIP (gather everything we accumulated)
        def build_zip():
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                # --- Summary text ---
                lines = ["FrostNode | Statistics Toolset — Auto Summary", ""]
                for name, blob in sorted(assets.get("texts", {}).items()):
                    text = blob.decode("utf-8", errors="replace") if isinstance(blob, (bytes, bytearray)) else str(blob)
                    lines.append(f"--- {name} ---")
                    lines.append(text)
                    lines.append("")  # blank line
                zf.writestr("summary.txt", "\n".join(lines))

                # --- Tables (CSV bytes or str) ---
                for name, blob in assets.get("tables", {}).items():
                    data = blob if isinstance(blob, (bytes, bytearray)) else str(blob).encode("utf-8")
                    zf.writestr(f"tables/{name}", data)

                # --- Figures (PNG bytes) ---
                for name, blob in assets.get("figs", {}).items():
                    data = blob if isinstance(blob, (bytes, bytearray)) else str(blob).encode("utf-8")
                    zf.writestr(f"figures/{name}", data)

            buf.seek(0)
            return buf.read()

        st.download_button("📦 Download-All (ZIP)", data=build_zip(), file_name="frostnode_export.zip", mime="application/zip")

        # Export to Excel workbook
        def export_xlsx():
            wb = Workbook()
            ws = wb.active; ws.title="Data"
            ws.append(list(df.columns))
            for _,row in df.iterrows():
                ws.append([row[c] for c in df.columns])
            # Summary sheet
            ws2 = wb.create_sheet("Summary")
            rowi = 1
            for k,v in assets["texts"].items():
                ws2.cell(row=rowi, column=1, value=k)
                ws2.cell(row=rowi+1, column=1, value=v.decode("utf-8"))
                rowi += 3
            # Tables sheet
            ws3 = wb.create_sheet("Tables")
            r = 1
            for name, blob in assets["tables"].items():
                ws3.cell(row=r, column=1, value=name); r+=1
                # write CSV-like content as plain text for simplicity
                ws3.cell(row=r, column=1, value=blob.decode("utf-8")); r+=2
            # Charts sheet (embed if possible)
            ws4 = wb.create_sheet("Charts")
            rimg = 1
            tmpdir = tempfile.mkdtemp()
            for name, blob in assets["figs"].items():
                p = os.path.join(tmpdir, name)
                with open(p, "wb") as f: f.write(blob)
                try:
                    img = XLImage(p)
                    ws4.add_image(img, f"A{rimg}")
                    rimg += 20
                except Exception:
                    pass
            # Formulas sheet (same as in app)
            ws5 = wb.create_sheet("Formulas")
            ws5.append(["Metric","Excel Formula (example ranges)"])
            rows = [
                ("Correlation (r)","=CORREL(B2:B101, A2:A101)"),
                ("R-squared","=RSQ(B2:B101, A2:A101)"),
                ("Slope (m)","=SLOPE(B2:B101, A2:A101)"),
                ("Intercept (b)","=INTERCEPT(B2:B101, A2:A101)"),
                ("ŷ at x0 (E2)","=FORECAST.LINEAR(E2, B2:B101, A2:A101)"),
                ("Std error (STEYX)","=STEYX(B2:B101, A2:A101)"),
                ("Mean (Y)","=AVERAGE(B2:B101)"),
                ("Median (Y)","=MEDIAN(B2:B101)"),
                ("Mode (Y)","=MODE.SNGL(B2:B101)"),
                ("Std Dev (Y sample)","=STDEV.S(B2:B101)"),
                ("Variance (Y sample)","=VAR.S(B2:B101)"),
                ("Q1/Q2/Q3 (Y)","=QUARTILE.EXC(B2:B101,1/2/3)"),
                ("Percentile 90% (Y)","=PERCENTILE.INC(B2:B101,0.9)"),
                ("Normal CDF","=NORM.S.DIST(z, TRUE)"),
                ("Normal INV","=NORM.S.INV(p)"),
            ]
            for a,b in rows: ws5.append([a,b])
            buf = io.BytesIO()
            wb.save(buf); buf.seek(0)
            return buf.read()

        st.download_button("⬇️ Export to Excel (.xlsx)", data=export_xlsx(), file_name="FrostNode_Export.xlsx",
                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        # Report: DOCX
        def export_docx():
            doc = Document()
            doc.add_heading("FrostNode | Statistics Toolset — Report", 0)
            for k,v in assets["texts"].items():
                doc.add_heading(k, level=2)
                doc.add_paragraph(v.decode("utf-8"))
            for name, blob in assets["figs"].items():
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                tmp.write(blob); tmp.flush()
                doc.add_picture(tmp.name, width=Inches(6))
                doc.paragraphs[-1].alignment = 1
                doc.add_paragraph(name.replace(".png",""))
            out = io.BytesIO()
            doc.save(out); out.seek(0)
            return out.read()

        st.download_button("📝 Export Word (.docx)", data=export_docx(), file_name="FrostNode_Report.docx",
                           mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

        # Report: PPTX
        def export_pptx():
            prs = Presentation()
            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "FrostNode | Statistics Toolset — Slides"
            slide.placeholders[1].text = "Auto-generated charts & captions"
            # Chart slides
            for name, blob in assets["figs"].items():
                blank = prs.slide_layouts[6]
                s = prs.slides.add_slide(blank)
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                tmp.write(blob); tmp.flush()
                pic = s.shapes.add_picture(tmp.name, PPTInches(1), PPTInches(1), width=PPTInches(8))
                tx = s.shapes.add_textbox(PPTInches(1), PPTInches(0.2), PPTInches(8), PPTInches(0.6))
                tx.text_frame.text = name.replace(".png","")
            out = io.BytesIO()
            prs.save(out); out.seek(0)
            return out.read()

        st.download_button("📽️ Export PowerPoint (.pptx)", data=export_pptx(), file_name="FrostNode_Charts.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# Footer
st.divider()
st.caption("Every chart has a Download PNG. Use the Exports tab to grab everything at once (ZIP, Excel, Word, PowerPoint).")
